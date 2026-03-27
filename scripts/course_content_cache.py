#!/usr/bin/env python3
"""Course Content Cache — dump course pages and files for Claude Code search.

Creates a local text cache of all Canvas pages and optionally extracts text from
uploaded course files (PDF, DOCX, PPTX). Claude Code can then grep/read the cache
directly for RAG-like Q&A without needing an embedding API.

Usage:
  python course_content_cache.py dump --course-id <id>
  python course_content_cache.py refresh --course-id <id>
  python course_content_cache.py status --course-id <id>
"""

import argparse
import json
import os
import re
import sys
import time
from pathlib import Path

# Logging
try:
    from idw_logger import get_logger
    _log = get_logger("course_content_cache")
    from idw_metrics import track as _track
except ImportError:
    import logging
    _log = logging.getLogger("course_content_cache")
    def _track(*a, **k): pass

# Canvas API — lazy import for testability
_canvas_api = None

def _get_canvas_api():
    global _canvas_api
    if _canvas_api is None:
        try:
            import canvas_api as _ca
            _canvas_api = _ca
        except ImportError:
            _log.error("canvas_api.py not found. Run from the scripts/ directory.")
            sys.exit(1)
    return _canvas_api

def get_config(**kwargs):
    return _get_canvas_api().get_config(**kwargs)

def request_with_retry(method, url, **kwargs):
    import requests as _req
    method_map = {"GET": _req.get, "POST": _req.post, "PUT": _req.put, "DELETE": _req.delete}
    callable_method = method_map.get(method, method) if isinstance(method, str) else method
    return _get_canvas_api()._request_with_retry(callable_method, url, **kwargs)

PLUGIN_ROOT = Path(__file__).resolve().parents[1]
CACHE_ROOT = PLUGIN_ROOT / "content_cache"


def _strip_html(html):
    """Strip HTML tags and normalize whitespace."""
    text = re.sub(r'<script[^>]*>.*?</script>', '', html, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r'<style[^>]*>.*?</style>', '', text, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r'<[^>]+>', ' ', text)
    text = re.sub(r'&nbsp;', ' ', text)
    text = re.sub(r'&amp;', '&', text)
    text = re.sub(r'&lt;', '<', text)
    text = re.sub(r'&gt;', '>', text)
    text = re.sub(r'&#\d+;', '', text)
    text = re.sub(r'\s+', ' ', text).strip()
    return text


def _cache_dir(course_id):
    """Get or create cache directory for a course."""
    d = CACHE_ROOT / str(course_id)
    (d / "pages").mkdir(parents=True, exist_ok=True)
    (d / "files").mkdir(parents=True, exist_ok=True)
    return d


def _load_metadata(course_id):
    """Load cache metadata (timestamps, counts)."""
    meta_path = _cache_dir(course_id) / "metadata.json"
    if meta_path.exists():
        return json.loads(meta_path.read_text())
    return {"course_id": str(course_id), "pages": {}, "files": {}, "last_dump": None}


def _save_metadata(course_id, metadata):
    """Save cache metadata."""
    meta_path = _cache_dir(course_id) / "metadata.json"
    meta_path.write_text(json.dumps(metadata, indent=2))


def dump_pages(course_id, since=None):
    """Fetch all Canvas pages and save as text files.

    Args:
        course_id: Canvas course ID
        since: Optional ISO timestamp — only fetch pages updated after this time

    Returns:
        Number of pages cached
    """
    config = get_config()
    if not config:
        raise RuntimeError("Canvas API not configured. Run /canvas-connect first.")

    cache = _cache_dir(course_id)
    metadata = _load_metadata(course_id)

    # Fetch all pages
    url = f"{config['base_url']}/courses/{course_id}/pages?per_page=100"
    if since:
        url += f"&sort=updated_at&order=desc"

    headers = {"Authorization": f"Bearer {config['token']}"}
    all_pages = []

    while url:
        resp = request_with_retry("GET", url, headers=headers)
        if resp.status_code != 200:
            _log.error("Failed to fetch pages: %s", resp.status_code)
            break
        all_pages.extend(resp.json())
        # Pagination
        link_header = resp.headers.get("Link", "")
        url = None
        if 'rel="next"' in link_header:
            for part in link_header.split(","):
                if 'rel="next"' in part:
                    url = part.split(";")[0].strip().strip("<>")
                    break

    count = 0
    for page in all_pages:
        slug = page.get("url", "")
        updated = page.get("updated_at", "")

        # Skip if not updated since last dump
        if since and updated and updated <= since:
            continue

        # Fetch full page body
        page_url = f"{config['base_url']}/courses/{course_id}/pages/{slug}"
        page_resp = request_with_retry("GET", page_url, headers=headers)
        if page_resp.status_code != 200:
            continue

        page_data = page_resp.json()
        body = page_data.get("body", "")
        title = page_data.get("title", slug)

        # Save HTML
        html_path = cache / "pages" / f"{slug}.html"
        html_path.write_text(body or "", encoding="utf-8")

        # Save plain text
        text = f"# {title}\n\n{_strip_html(body)}"
        txt_path = cache / "pages" / f"{slug}.txt"
        txt_path.write_text(text, encoding="utf-8")

        metadata["pages"][slug] = {
            "title": title,
            "updated_at": updated,
            "cached_at": time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime()),
            "size": len(body),
        }
        count += 1
        _log.info("Cached: %s (%s)", title, slug)

    metadata["last_dump"] = time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime())
    _save_metadata(course_id, metadata)
    _track("content_cached", count)
    return count


def dump_files(course_id):
    """Download course files and extract text where possible.

    Supports PDF (pymupdf), DOCX (python-docx), PPTX (python-pptx), TXT, and MD.
    """
    config = get_config()
    if not config:
        raise RuntimeError("Canvas API not configured.")

    import requests as req_lib

    cache = _cache_dir(course_id)
    metadata = _load_metadata(course_id)

    url = f"{config['base_url']}/courses/{course_id}/files?per_page=100"
    headers = {"Authorization": f"Bearer {config['token']}"}
    resp = request_with_retry("GET", url, headers=headers)
    if resp.status_code != 200:
        _log.error("Failed to fetch files: %s", resp.status_code)
        return 0

    files = resp.json()
    count = 0

    for file_info in files:
        name = file_info.get("display_name", "")
        file_url = file_info.get("url", "")
        ext = os.path.splitext(name)[1].lower()

        # Only process text-extractable files
        if ext not in (".pdf", ".docx", ".pptx", ".txt", ".md", ".html", ".htm"):
            continue

        try:
            _log.info("Downloading: %s", name)
            dl_resp = req_lib.get(file_url, timeout=30)
            if dl_resp.status_code != 200:
                continue

            text = ""
            if ext in (".txt", ".md"):
                text = dl_resp.content.decode("utf-8", errors="ignore")
            elif ext in (".html", ".htm"):
                text = _strip_html(dl_resp.content.decode("utf-8", errors="ignore"))
            elif ext == ".pdf":
                try:
                    import fitz  # pymupdf
                    doc = fitz.open(stream=dl_resp.content, filetype="pdf")
                    text = "\n\n".join(page.get_text() for page in doc)
                    doc.close()
                except ImportError:
                    _log.warning("pymupdf not installed — skipping PDF: %s", name)
                    continue
            elif ext == ".docx":
                try:
                    import docx
                    import io
                    doc = docx.Document(io.BytesIO(dl_resp.content))
                    text = "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
                except ImportError:
                    _log.warning("python-docx not installed — skipping DOCX: %s", name)
                    continue
            elif ext == ".pptx":
                try:
                    from pptx import Presentation
                    import io
                    prs = Presentation(io.BytesIO(dl_resp.content))
                    slides_text = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                slides_text.append(shape.text_frame.text)
                    text = "\n\n".join(slides_text)
                except ImportError:
                    _log.warning("python-pptx not installed — skipping PPTX: %s", name)
                    continue

            if text.strip():
                safe_name = re.sub(r'[^\w\-.]', '_', name)
                txt_path = cache / "files" / f"{safe_name}.txt"
                txt_path.write_text(f"# {name}\n\n{text}", encoding="utf-8")
                metadata["files"][name] = {
                    "cached_at": time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime()),
                    "size": len(text),
                }
                count += 1
                _log.info("Extracted text: %s (%d chars)", name, len(text))

        except Exception as e:
            _log.warning("Error processing %s: %s", name, e)

    _save_metadata(course_id, metadata)
    return count


def get_status(course_id):
    """Get cache status for a course."""
    metadata = _load_metadata(course_id)
    cache = _cache_dir(course_id)

    pages_dir = cache / "pages"
    files_dir = cache / "files"
    page_count = len(list(pages_dir.glob("*.txt"))) if pages_dir.exists() else 0
    file_count = len(list(files_dir.glob("*.txt"))) if files_dir.exists() else 0

    return {
        "course_id": str(course_id),
        "page_count": page_count,
        "file_count": file_count,
        "last_dump": metadata.get("last_dump"),
        "cache_path": str(cache),
    }


# ── CLI ─────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(description="Course content cache for Claude Code search")
    parser.add_argument("command", choices=["dump", "refresh", "status"],
                        help="dump: full cache, refresh: only updated pages, status: show cache info")
    parser.add_argument("--course-id", required=True, help="Canvas course ID")
    parser.add_argument("--include-files", action="store_true",
                        help="Also download and extract text from course files")
    args = parser.parse_args()
    _track("skill_invoked", context={"skill": "knowledge"})

    if args.command == "status":
        status = get_status(args.course_id)
        print(json.dumps(status, indent=2))
        return

    if args.command == "dump":
        count = dump_pages(args.course_id)
        print(f"Cached {count} pages")
        if args.include_files:
            file_count = dump_files(args.course_id)
            print(f"Cached {file_count} files")

    elif args.command == "refresh":
        metadata = _load_metadata(args.course_id)
        since = metadata.get("last_dump")
        count = dump_pages(args.course_id, since=since)
        print(f"Refreshed {count} pages (since {since or 'never'})")

    status = get_status(args.course_id)
    print(json.dumps(status, indent=2))


if __name__ == "__main__":
    try:
        main()
    except KeyboardInterrupt:
        sys.exit(130)
    except Exception as e:
        _log.exception("Unexpected error")
        print(f"\nSomething went wrong: {e}")
        print("Check the log at ~/.idw/logs/ for details.")
        sys.exit(1)
